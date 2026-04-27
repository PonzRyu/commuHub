from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PATH = ROOT / "PORTFOLIO.pptx"


def apply_title_style(shape) -> None:
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.font.name = "Meiryo"
    paragraph.font.size = Pt(36)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(15, 23, 42)


def apply_subtitle_style(shape) -> None:
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.font.name = "Meiryo"
    paragraph.font.size = Pt(18)
    paragraph.font.color.rgb = RGBColor(71, 85, 105)


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "CommuHub"
    slide.placeholders[1].text = (
        "フロントエンドエンジニア応募用ポートフォリオ\n"
        "週間スケジュール共有Webアプリ"
    )
    apply_title_style(slide.shapes.title)
    apply_subtitle_style(slide.placeholders[1])


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    title_paragraph = slide.shapes.title.text_frame.paragraphs[0]
    title_paragraph.font.name = "Meiryo"
    title_paragraph.font.size = Pt(30)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(15, 23, 42)

    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    for index, bullet in enumerate(bullets):
        paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.name = "Meiryo"
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(30, 41, 59)
        paragraph.space_after = Pt(8)
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT


def add_two_column_slide(
    prs: Presentation,
    title: str,
    left_title: str,
    left_bullets: list[str],
    right_title: str,
    right_bullets: list[str],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    title_paragraph = slide.shapes.title.text_frame.paragraphs[0]
    title_paragraph.font.name = "Meiryo"
    title_paragraph.font.size = Pt(30)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(15, 23, 42)

    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.8), Inches(4.8))
    right_box = slide.shapes.add_textbox(Inches(6.7), Inches(1.6), Inches(5.8), Inches(4.8))

    for box, section_title, section_bullets in (
        (left_box, left_title, left_bullets),
        (right_box, right_title, right_bullets),
    ):
        tf = box.text_frame
        tf.clear()
        header = tf.paragraphs[0]
        header.text = section_title
        header.font.name = "Meiryo"
        header.font.size = Pt(22)
        header.font.bold = True
        header.font.color.rgb = RGBColor(15, 23, 42)
        header.space_after = Pt(10)

        for item in section_bullets:
            paragraph = tf.add_paragraph()
            paragraph.text = f"・{item}"
            paragraph.level = 0
            paragraph.font.name = "Meiryo"
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = RGBColor(30, 41, 59)
            paragraph.space_after = Pt(6)


def build_presentation() -> Presentation:
    prs = Presentation()

    add_title_slide(prs)

    add_bullet_slide(
        prs,
        "プロジェクト概要",
        [
            "チーム内予定表でメンバーの週間日程を可視化",
            "個別カレンダー確認の手間を削減し、会議調整を効率化",
            "誰が・いつ・何の予定かを1画面で把握できるUIを設計",
        ],
    )

    add_bullet_slide(
        prs,
        "担当範囲",
        [
            "週間スケジュール画面の設計・実装",
            "管理画面（部署・メンバー・ICS管理）の実装",
            "週次共有事項・週次予定表機能の実装",
            "バリデーション、エラーハンドリング、セッション管理の実装",
        ],
    )

    add_two_column_slide(
        prs,
        "主要機能",
        "ユーザー向け機能",
        [
            "週間スケジュール一覧",
            "部署フィルタ、週移動",
            "予定の手動リフレッシュ",
            "予定なしメンバーの非表示切替",
        ],
        "管理者向け機能",
        [
            "部署CRUD",
            "メンバーCRUD",
            "ICSリンク登録・更新・削除",
            "アプリ表示名、外部リンク管理",
        ],
    )

    add_bullet_slide(
        prs,
        "フロントエンド実装の工夫",
        [
            "Next.js App RouterでServer/Client責務を分離",
            "dnd-kit導入時のハイドレーション不一致を回避",
            "モーダル中心の編集UXで操作の文脈を維持",
            "入力検証とエラー表示で再入力しやすい設計",
        ],
    )

    add_two_column_slide(
        prs,
        "技術スタック",
        "Frontend / UI",
        [
            "Next.js (App Router), React, TypeScript",
            "Tailwind CSS, shadcn/ui, Radix UI",
            "dnd-kit",
        ],
        "Backend / Data",
        [
            "Next.js Server Actions",
            "Prisma + PostgreSQL",
            "node-ical (ICS展開)",
            "Electron (Windows配布対応)",
        ],
    )

    add_bullet_slide(
        prs,
        "セキュリティと品質",
        [
            "管理画面へのアクセスをセッションで制御",
            "遷移先パスを安全化し不正リダイレクトを抑止",
            "ICS URLの入力検証（protocol/localhost/private IP）",
            "運用を意識したデータ検証とエラーハンドリング",
        ],
    )

    add_bullet_slide(
        prs,
        "今後の改善",
        [
            "E2Eテスト導入（Playwright）",
            "大規模データ時の表示パフォーマンス最適化",
            "監査ログ・操作履歴の追加",
            "権限レベルの細分化",
        ],
    )

    add_bullet_slide(
        prs,
        "応募時メッセージ",
        [
            "実運用を意識し、一覧性・操作性・保守性を重視して実装",
            "フロントエンドでの体験設計と堅牢な入力検証を両立",
            "要件に沿って機能を組み上げるだけでなく、運用性も考慮",
        ],
    )

    return prs


def main() -> None:
    presentation = build_presentation()
    presentation.save(OUTPUT_PATH)
    print(f"Generated: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
